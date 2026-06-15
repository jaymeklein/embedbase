"""Generate the binary fixtures for the docling parser tests.

These fixtures are committed to the repo so the gated tests in
``tests/unit/test_parsers_docling.py`` can run against the *real* docling library
without regenerating binaries. Regenerate them only if the expected structure
needs to change.

The generator deliberately lives outside the normal dependency set: it needs
``reportlab`` (tabled PDF), ``Pillow`` (image-only scanned PDF), ``python-docx``
and ``python-pptx`` (office formats). ``Pillow``/``python-docx``/``python-pptx``
ship with docling; ``reportlab`` must be installed separately.

Run from the repo root inside an environment that has docling installed::

    python tests/fixtures/docling/generate_fixtures.py
"""

from __future__ import annotations

from pathlib import Path

_OUT = Path(__file__).resolve().parent


def make_scanned_pdf(path: Path) -> None:
    """Write a 2-page *image-only* PDF (no text layer -> OCR required)."""
    from PIL import Image, ImageDraw, ImageFont

    try:
        font = ImageFont.truetype("arial.ttf", 64)
    except OSError:
        font = ImageFont.load_default(64)

    pages: list[Image.Image] = []
    for text in (
        "INVOICE NUMBER 12345\nTotal Due 980 USD",
        "Thank you for your business\nPayment terms net 30 days",
    ):
        img = Image.new("RGB", (1240, 1754), "white")  # ~A4 at 150 DPI
        ImageDraw.Draw(img).multiline_text((120, 240), text, fill="black", font=font, spacing=24)
        pages.append(img)
    pages[0].save(path, save_all=True, append_images=pages[1:], format="PDF", resolution=150.0)


def make_table_pdf(path: Path) -> None:
    """Write a single-page PDF: heading + prose + a ruled (grid) table + prose.

    The surrounding prose matters: docling's layout model classifies a near-empty
    page as a single picture, so the table needs to sit inside a realistic document
    body for the table region to be recognised and serialised as Markdown.
    """
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

    data = [
        ["Region", "Q1 Revenue", "Q2 Revenue", "Growth"],
        ["North", "100", "120", "20%"],
        ["South", "90", "115", "28%"],
        ["West", "70", "88", "26%"],
        ["East", "140", "150", "7%"],
        ["Central", "60", "95", "58%"],
    ]
    table = Table(data, colWidths=[120, 110, 110, 80])
    table.setStyle(
        TableStyle(
            [
                ("GRID", (0, 0), (-1, -1), 1, colors.black),
                ("BACKGROUND", (0, 0), (-1, 0), colors.lightgrey),
                ("FONTSIZE", (0, 0), (-1, -1), 12),
                ("TOPPADDING", (0, 0), (-1, -1), 8),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
            ]
        )
    )
    styles = getSampleStyleSheet()
    intro = (
        "The following table summarises quarterly sales by region. Each row lists a "
        "region together with its revenue for the first two quarters of the year."
    )
    outro = (
        "Central showed the strongest growth, while East remained the largest region "
        "by absolute revenue across both quarters of the reporting period."
    )
    SimpleDocTemplate(str(path), pagesize=LETTER).build(
        [
            Paragraph("Quarterly Sales Report", styles["Title"]),
            Spacer(1, 16),
            Paragraph(intro, styles["BodyText"]),
            Spacer(1, 16),
            table,
            Spacer(1, 16),
            Paragraph(outro, styles["BodyText"]),
        ]
    )


def make_docx(path: Path) -> None:
    """Write a .docx with three headings, each followed by a body paragraph."""
    from docx import Document

    doc = Document()
    for heading, body in (
        ("Introduction", "This document introduces the embedding system and its goals."),
        ("Architecture", "The architecture separates the API service from the worker."),
        ("Deployment", "Deployment uses Docker Compose with optional GPU acceleration."),
    ):
        doc.add_heading(heading, level=1)
        doc.add_paragraph(body)
    doc.save(str(path))


def make_pptx(path: Path) -> None:
    """Write a .pptx with two slides, each carrying a title + body text."""
    from pptx import Presentation

    prs = Presentation()
    layout = prs.slide_layouts[1]  # "Title and Content"
    for title, body in (
        ("Overview", "EmbedBase indexes documents into vector embeddings."),
        ("Roadmap", "Delivery 4 adds the MCP server and the docling parser."),
    ):
        slide = prs.slides.add_slide(layout)
        title_shape = slide.shapes.title
        body_shape = slide.placeholders[1]
        title_shape.text = title
        body_shape.text = body
    prs.save(str(path))


def main() -> None:
    """Generate every fixture into this directory."""
    make_scanned_pdf(_OUT / "scanned_2page.pdf")
    make_table_pdf(_OUT / "table.pdf")
    make_docx(_OUT / "headings.docx")
    make_pptx(_OUT / "slides.pptx")
    print(f"Wrote fixtures to {_OUT}")


if __name__ == "__main__":
    main()
